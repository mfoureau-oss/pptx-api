from flask import Flask, request, send_file
from pptx import Presentation
from pptx.util import Inches
import requests
import io
import tempfile
import os

app = Flask(__name__)

@app.route("/update-pptx", methods=["POST"])
def update_pptx():
    data = request.json
    slide1_text = data.get("slide1_text", "")
    slide4_images = data.get("slide4_images", [])  # list of image URLs
    pptx_url = data.get("pptx_url")

    if not pptx_url:
        return {"error": "pptx_url is required"}, 400

    # Download the PPTX file
    pptx_response = requests.get(pptx_url)
    if pptx_response.status_code != 200:
        return {"error": "Unable to download PPTX"}, 500

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.pptx")
        output_path = os.path.join(tmpdir, "output.pptx")

        with open(input_path, "wb") as f:
            f.write(pptx_response.content)

        prs = Presentation(input_path)

        # Modify slide 1
        try:
            slide1 = prs.slides[0]
            for shape in slide1.shapes:
                if shape.has_text_frame:
                    shape.text = slide1_text
                    break
        except Exception as e:
            return {"error": f"Error editing slide 1: {str(e)}"}, 500

        # Add images to slide 4
        try:
            slide4 = prs.slides[3]
            left = Inches(0.5)
            top = Inches(1)
            width = Inches(4.5)
            for idx, img_url in enumerate(slide4_images[:3]):
                img_data = requests.get(img_url).content
                image_stream = io.BytesIO(img_data)
                slide4.shapes.add_picture(image_stream, left, top + Inches(idx * 2.5), width=width)
        except Exception as e:
            return {"error": f"Error editing slide 4: {str(e)}"}, 500

        # Save modified presentation
        prs.save(output_path)

        return send_file(output_path, as_attachment=True, download_name="modified_presentation.pptx")

if __name__ == "__main__":
    app.run(debug=True, host="0.0.0.0", port=10000)
